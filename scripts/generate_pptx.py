import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths to assets
ARTIFACT_DIR = "/home/michael/.gemini/antigravity/brain/87f2f969-e14e-4789-b5a8-933490eb77b7"
DOCS_DIR = "/home/michael/proj/dcgame/docs/marketing"
OUTPUT_DIR = "/home/michael/proj/dcgame/marketing/presentations"

# Asset Pack Mapping
ASSET_PACK = {
    "FEAR": os.path.join(ARTIFACT_DIR, "demon_portrait_fear_1772179292929.png"),
    "DOUBT": os.path.join(ARTIFACT_DIR, "demon_portrait_doubt_1772179311024.png"),
    "SWORD": os.path.join(ARTIFACT_DIR, "powerup_sword_spirit_1772179324845.png"),
    "ARMOR": os.path.join(ARTIFACT_DIR, "powerup_shield_faith_helmet_salvation_1772179338041.png"),
    "QUIZ": os.path.join(ARTIFACT_DIR, "quiz_mode_icons_set_1772179419552.png"),
    "MAPS": os.path.join(ARTIFACT_DIR, "map_style_grid_city_labyrinth_1772179434148.png"),
}

TITLE_BG_MAP = {
    "player-angle-1-fun-adventure.md": "versebattles_title_background_1772171807202.png",
    "player-angle-2-learn-grow.md": "versebattles_learn_grow_theme_1772173221717.png",
    "player-angle-3-relax-meditate.md": "versebattles_relax_meditate_theme_1772173236387.png",
    "player-angle-4-challenge-progress.md": "versebattles_challenge_progress_theme_1772173249736.png",
    "player-angle-5-family-community.md": "versebattles_family_community_theme_1772173263242.png",
}

LANDING_PAGE = os.path.join(ARTIFACT_DIR, "landing_page_1772171194901.png")
GAMEPLAY_ACTION = os.path.join(ARTIFACT_DIR, "gameplay_action_1772171657859.png")

def parse_md(filepath):
    """Simple parser to extract sections from the MD files."""
    with open(filepath, 'r') as f:
        content = f.read()
    
    # Ensure branding is consistent
    content = content.replace("dcgame.4you.tel", "Versebattles.com")
    
    sections = re.split(r'\n---\n', content)
    parsed_sections = []
    
    for section in sections:
        lines = section.strip().split('\n')
        if not lines: continue
        title = ""
        body_lines = []
        for line in lines:
            if line.startswith('#'):
                title = line.replace('#', '').strip()
            elif line.startswith('###'):
                title = line.replace('###', '').strip()
            else:
                clean_line = re.sub(r'[*_!\[\]()]', '', line).strip()
                if clean_line:
                    body_lines.append(clean_line)
        
        parsed_sections.append({"title": title, "body_lines": body_lines})
    
    return parsed_sections

def add_content_slide(prs, title, body_lines, image_path=None):
    """Adds a slide, possibly splitting content if it's too long."""
    MAX_LINES = 7
    if image_path:
        MAX_LINES = 5
    
    chunks = [body_lines[i:i + MAX_LINES] for i in range(0, len(body_lines), MAX_LINES)]
    
    for idx, chunk in enumerate(chunks):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title if idx == 0 else f"{title} (Cont.)"
        
        content_shape = slide.placeholders[1]
        
        # Smart positioning if image exists
        if image_path and idx == 0:
            content_shape.width = Inches(4.5)
            slide.shapes.add_picture(image_path, Inches(5.2), Inches(2), width=Inches(4))
        
        content_shape.text = "\n".join(chunk)

def create_presentation(md_filename):
    md_path = os.path.join(DOCS_DIR, md_filename)
    bg_image = os.path.join(ARTIFACT_DIR, TITLE_BG_MAP[md_filename])
    output_name = md_filename.replace('.md', '.pptx')
    output_path = os.path.join(OUTPUT_DIR, output_name)
    
    sections = parse_md(md_path)
    prs = Presentation()

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(bg_image, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    if sections:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.text = sections[0]["title"]
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 2. Iterate and Inject Assets Based on Keyword Matching
    for i, section in enumerate(sections[1:]):
        text_blob = " ".join(section["body_lines"]).lower()
        title_blob = section["title"].lower()
        
        # Keyword-based image injection
        injected_image = None
        if "demon" in text_blob or "fear" in text_blob:
            injected_image = ASSET_PACK["FEAR"]
        elif "doubt" in title_blob or "doubt" in text_blob:
            injected_image = ASSET_PACK["DOUBT"]
        elif "sword" in text_blob or "combat" in text_blob:
            injected_image = ASSET_PACK["SWORD"]
        elif "shield" in text_blob or "armor" in text_blob:
            injected_image = ASSET_PACK["ARMOR"]
        elif "quiz" in text_blob or "mode" in text_blob:
            injected_image = ASSET_PACK["QUIZ"]
        elif "map" in text_blob or "labyrinth" in text_blob:
            injected_image = ASSET_PACK["MAPS"]
        elif i == 1:
            injected_image = GAMEPLAY_ACTION
        elif i == len(sections) - 2:
            injected_image = LANDING_PAGE
            
        add_content_slide(prs, section["title"], section["body_lines"], injected_image)

    prs.save(output_path)
    print(f"Enhanced: {output_path}")

if __name__ == "__main__":
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
    for md_file in TITLE_BG_MAP.keys():
        create_presentation(md_file)
