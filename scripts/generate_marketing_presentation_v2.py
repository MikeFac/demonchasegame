from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/michael/proj/dcgame")
OUTPUT = ROOT / "docs/marketing/presentations/VerseBattles_Marketing_v2.pptx"
CHAPTERS_JSON = ROOT / "missions/chapters.json"

LOGO = ROOT / "images/VerseBattles-logo.png"
TITLE_BG = ROOT / "images/splash_hopeful_clean.png"

IMAGES = {
    "menu": ROOT / "public/landing/menu-screen.png",
    "gameplay": ROOT / "public/landing/gameplay-screen.png",
    "worlds": ROOT / "public/landing/worlds-screen.png",
    "tutorial": ROOT / "public/landing/tutorial-screen.png",
    "start_here": ROOT / "output/web-game/start-here-summary-final/summary-end-to-end.png",
    "review": ROOT / "output/web-game/review-toast-leak/review-no-toast.png",
    "multiplayer": ROOT / "output/web-game/multiplayer-regression/pageA-final.png",
    "wave": ROOT / "output/web-game/wave-assault-smoke/wave-initial.png",
    "maze": ROOT / "output/web-game/scripture-maze-direct/shot-0.png",
}

BG = RGBColor(13, 18, 28)
PANEL = RGBColor(22, 31, 48)
PANEL_ALT = RGBColor(29, 43, 66)
TEXT = RGBColor(240, 244, 250)
MUTED = RGBColor(172, 188, 210)
ACCENT = RGBColor(115, 220, 171)
ACCENT_2 = RGBColor(108, 165, 255)
WARN = RGBColor(255, 208, 108)


def mission_stats() -> tuple[int, int]:
    with CHAPTERS_JSON.open() as f:
        data = json.load(f)
    chapters = data["chapters"]
    mission_total = sum(len(chapter.get("missionIds", [])) for chapter in chapters)
    return len(chapters), mission_total


def add_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, title: str, kicker: str | None = None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.33),
        Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PANEL
    bar.line.color.rgb = PANEL

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(8.7), Inches(0.4))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(25)
    run.font.bold = True
    run.font.color.rgb = TEXT

    if kicker:
        kicker_box = slide.shapes.add_textbox(Inches(9.45), Inches(0.2), Inches(3.2), Inches(0.35))
        tf = kicker_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = kicker
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = ACCENT


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.1), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT


def add_panel(slide, left: float, top: float, width: float, height: float, color=PANEL_ALT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_image(slide, image_path: Path, left: float, top: float, width: float, height: float | None = None):
    if not image_path.exists():
        return None
    if height is None:
        return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_caption(slide, title: str, body: str, left: float, top: float, width: float, height: float):
    panel = add_panel(slide, left, top, width, height, color=PANEL)
    panel.fill.transparency = 0.1
    box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.16), Inches(width - 0.36), Inches(height - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = body
    p2.space_before = Pt(3)
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT


def add_title_slide(prs: Presentation, chapters: int, missions: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_image(slide, TITLE_BG, 0, 0, 13.33, 7.5)

    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG
    overlay.fill.transparency = 0.35
    overlay.line.color.rgb = BG

    add_image(slide, LOGO, 0.75, 0.45, 1.3)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(8.5), Inches(2.1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VerseBattles"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Marketing Presentation V2"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = "Updated for the current product: missions, guided onboarding, discipleship tracks, verse-song learning, multiplayer, and ministry-facing pilots."
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT

    stat_panel = add_panel(slide, 0.8, 4.8, 5.7, 1.35)
    stat_panel.fill.transparency = 0.12
    stat_box = slide.shapes.add_textbox(Inches(1.05), Inches(5.05), Inches(5.2), Inches(0.85))
    tf = stat_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{chapters} chapters · {missions} missions · 5 quiz modes · solo + multiplayer + review + Verse of the Day"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT

    add_footer(slide, f"Generated from current repo state · {date.today().isoformat()}")


def add_standard_slide(prs: Presentation, title: str, bullets: list[str], image_key: str, kicker: str | None = None, caption: tuple[str, str] | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_bar(slide, title, kicker=kicker)
    add_bullets(slide, bullets, 0.7, 1.25, 5.0, 5.4)
    add_panel(slide, 5.85, 1.18, 6.75, 5.55)
    add_image(slide, IMAGES[image_key], 6.02, 1.34, 6.42, 4.8)
    if caption:
        add_caption(slide, caption[0], caption[1], 6.08, 6.0, 6.2, 0.5)
    add_footer(slide, "versebattles.com")


def add_modes_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_bar(slide, "Current Play Surface", kicker="What people can use now")

    left_bullets = [
        "Solo combat runs for quick individual play.",
        "Missions provide structured progression instead of only endless runs.",
        "Multiplayer supports shared sessions rather than isolated solo use.",
        "Review / Learn and Verse of the Day extend the product beyond combat.",
        "The menu also exposes a 3D experimental view for alternate presentation.",
    ]
    add_bullets(slide, left_bullets, 0.7, 1.25, 4.9, 5.4)

    add_panel(slide, 5.8, 1.16, 3.15, 2.55)
    add_panel(slide, 9.05, 1.16, 3.15, 2.55)
    add_panel(slide, 5.8, 3.86, 6.4, 2.55)
    add_image(slide, IMAGES["worlds"], 5.95, 1.3, 2.85, 2.15)
    add_image(slide, IMAGES["multiplayer"], 9.2, 1.3, 2.85, 2.15)
    add_image(slide, IMAGES["review"], 5.95, 4.0, 6.1, 2.1)

    add_caption(slide, "Missions / menu", "Structured navigation and multiple game paths are already visible from the product surface.", 5.95, 3.18, 2.7, 0.45)
    add_caption(slide, "Multiplayer", "Shared play is already part of the real product, not just a roadmap promise.", 9.2, 3.18, 2.7, 0.45)
    add_caption(slide, "Review and learning", "Players can move from combat into verse review, songs, and devotional follow-up.", 6.05, 6.18, 5.9, 0.38)
    add_footer(slide, "versebattles.com")


def add_missions_slide(prs: Presentation, chapters: int, missions: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_bar(slide, "Missions, Onboarding, and Discipleship Tracks", kicker=f"{chapters} chapters · {missions} missions")
    bullets = [
        "A guided Start Here first win now helps new players grasp the loop quickly.",
        "The mission map is no longer a tiny side feature; it is a real progression path.",
        "The Teachings of Jesus chapter turns the game toward discipleship content, not only arcade pressure.",
        "Wave Assault and Scripture Maze add alternate mission styles inside the same product.",
        "This makes the product easier to pilot with groups because there is more than one way to enter it.",
    ]
    add_bullets(slide, bullets, 0.7, 1.22, 5.05, 5.55)

    add_panel(slide, 5.92, 1.16, 6.58, 5.55)
    add_image(slide, IMAGES["start_here"], 6.08, 1.34, 3.12, 2.25)
    add_image(slide, IMAGES["wave"], 9.35, 1.34, 2.95, 2.25)
    add_image(slide, IMAGES["maze"], 6.08, 3.82, 2.95, 2.45)
    add_image(slide, IMAGES["tutorial"], 9.18, 3.82, 3.12, 2.45)

    add_caption(slide, "Start Here", "A guided first-run mission now exists instead of dropping new players into the deep end.", 6.1, 3.35, 2.95, 0.38)
    add_caption(slide, "Wave Assault", "Alternate mission format for more arcade-style sessions.", 9.38, 3.35, 2.72, 0.38)
    add_caption(slide, "Scripture Maze", "Another playable mission variant, not just a design doc.", 6.1, 6.12, 2.72, 0.38)
    add_caption(slide, "Landing/tutorial surface", "Screenshots and audience pages make it easier to explain the product to adults before they play.", 9.18, 6.12, 3.05, 0.38)
    add_footer(slide, "versebattles.com")


def add_discipleship_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_bar(slide, "Deeper Scripture Follow-Up", kicker="Current learning features")
    bullets = [
        "Review mode gives players a place to slow down and study after action gameplay.",
        "Verse of the Day adds a lighter-weight repeatable learning entry point.",
        "AI devotional sermon / prayer screens add reflective follow-up around a verse.",
        "The integrated verse-song library gives another retention path beyond text recall.",
        "This is materially different from an arcade shell with Bible trivia pasted on top.",
    ]
    add_bullets(slide, bullets, 0.7, 1.22, 5.0, 5.55)
    add_panel(slide, 5.88, 1.16, 6.6, 5.55)
    add_image(slide, IMAGES["review"], 6.05, 1.34, 6.25, 4.8)
    add_caption(slide, "Learning stack", "Review, devotionals, prayer, and songs now give the product more discipleship depth than the original marketing deck showed.", 6.1, 6.02, 6.05, 0.46)
    add_footer(slide, "versebattles.com")


def add_positioning_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_bar(slide, "Who This Fits Best Right Now", kicker="Current positioning")

    add_panel(slide, 0.7, 1.28, 3.9, 4.7)
    add_panel(slide, 4.72, 1.28, 3.9, 4.7)
    add_panel(slide, 8.74, 1.28, 3.9, 4.7)

    columns = [
        ("Youth pastors", [
            "Use as a warm-up, challenge, or pilot discipleship tool.",
            "Strongest near-term distribution wedge because one leader can bring a group.",
            "Dedicated landing page already exists: /youth-pastors",
        ]),
        ("Missions leaders", [
            "Useful where browser access and low setup matter.",
            "Missions-facing positioning already exists: /missions",
            "Better framed as a practical pilot than a mass-market app.",
        ]),
        ("Parents", [
            "A spiritually meaningful alternative to low-value screen time.",
            "Works on household devices without app-store dependency.",
            "Dedicated landing page already exists: /parents",
        ]),
    ]

    x_positions = [0.95, 4.97, 8.99]
    for (title, bullets), x in zip(columns, x_positions):
        box = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(3.35), Inches(4.05))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.bullet = True
            p.space_before = Pt(8)
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT

    note = slide.shapes.add_textbox(Inches(0.88), Inches(6.2), Inches(11.6), Inches(0.45))
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The v2 story is not \"look at a concept.\" It is \"here is a usable product with clearer audience fits and better pilotability.\""
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = WARN
    add_footer(slide, "versebattles.com")


def add_closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_image(slide, TITLE_BG, 0, 0, 13.33, 7.5)

    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG
    overlay.fill.transparency = 0.22
    overlay.line.color.rgb = BG

    add_image(slide, LOGO, 0.8, 0.55, 1.15)

    panel = add_panel(slide, 0.8, 1.3, 6.35, 4.8)
    panel.fill.transparency = 0.1
    box = slide.shapes.add_textbox(Inches(1.08), Inches(1.6), Inches(5.7), Inches(4.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Current ask"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    for line in [
        "Test the actual game at versebattles.com.",
        "Run a small pilot with a youth group, family, or ministry setting.",
        "Report where the product helps and where it still breaks down.",
        "If a specific ministry use case is strong, sponsor that rollout rather than funding vague platform hype.",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    box2 = slide.shapes.add_textbox(Inches(7.6), Inches(1.75), Inches(4.7), Inches(3.5))
    tf = box2.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "What this deck does differently"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TEXT

    for line in [
        "Leads with current capabilities, not only future vision.",
        "Reflects missions, onboarding, discipleship tracks, and learning layers already in the repo.",
        "Keeps future concepts separate from the shipped product.",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT

    add_footer(slide, "VerseBattles Marketing Presentation V2")


def build_deck():
    chapters, missions = mission_stats()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, chapters, missions)
    add_standard_slide(
        prs,
        "What VerseBattles Is Now",
        [
            "A playable browser-based Scripture game, not just a product idea.",
            "Players answer Bible-verse prompts to earn attacks, survive pressure, and progress.",
            "The product now includes multiple learning paths instead of a single arcade loop.",
            "It works as a game first, while still giving adults a serious discipleship angle.",
            "No app-store install is required to test it on phones or desktop.",
        ],
        "menu",
        kicker="Current-state summary",
        caption=("Live product surface", "The menu already exposes missions, learning, multiplayer, and configuration rather than a single narrow entry point."),
    )
    add_standard_slide(
        prs,
        "Core Gameplay Loop",
        [
            "Move, dodge, and survive real combat pressure.",
            "Answer Scripture prompts correctly to gain ammo and deal damage.",
            "Five quiz modes create recall variety: first-letter, missing-word, category-match, true/false, and cloze.",
            "Verse categories and combat choices matter alongside reflexes.",
            "This gives the product a stronger game feel than a static flashcard wrapper.",
        ],
        "gameplay",
        kicker="The hook",
        caption=("Gameplay capture", "The moment-to-moment loop is still the product's attention engine."),
    )
    add_modes_slide(prs)
    add_missions_slide(prs, chapters, missions)
    add_discipleship_slide(prs)
    add_standard_slide(
        prs,
        "Why It Is More Pilot-Ready Than Before",
        [
            "The guided Start Here flow reduces first-session confusion.",
            "Missions give adults a clearer structure to test than open-ended solo play alone.",
            "Audience-specific landing pages now exist for youth pastors, parents, and missions leaders.",
            "Learning follow-up is stronger because review, devotionals, and songs are now present.",
            "This supports a more credible ask: run a pilot and judge the real product.",
        ],
        "tutorial",
        kicker="Why v2 exists",
        caption=("Acquisition and onboarding", "The product and its marketing surfaces now tell a more coherent story."),
    )
    add_positioning_slide(prs)
    add_closing_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build_deck()
