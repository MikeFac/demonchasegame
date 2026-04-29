from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/michael/proj/dcgame")
OUT = ROOT / "docs/marketing/presentations/VerseBattles_Uganda_Christian_Schools.pptx"

LOGO = ROOT / "images/VerseBattles-logo.png"
SPLASH_HOPEFUL = ROOT / "images/splash_hopeful_clean.png"
SPLASH_STRIKE = ROOT / "images/splash_strike.png"
MENU_SCREEN = ROOT / "public/landing/menu-screen.png"
GAMEPLAY = ROOT / "public/landing/gameplay-screen.png"
MISSIONS = ROOT / "output/web-game/onboarding-missions/shot-0.png"
REVIEW = ROOT / "output/web-game/review-toast-leak/review-no-toast.png"
MULTI = ROOT / "output/web-game/multiplayer-regression/pageA-final.png"

MENU_ICONS = [
    ROOT / "images/menu/solo_game_labeled_icon.png",
    ROOT / "images/menu/missions_labeled_icon.png",
    ROOT / "images/menu/learn_verses_labeled_icon.png",
    ROOT / "images/menu/multiplayer_labeled_icon.png",
    ROOT / "images/menu/groups_labeled_icon.png",
]

BG = RGBColor(14, 20, 31)
PANEL = RGBColor(24, 34, 52)
PANEL_2 = RGBColor(33, 49, 77)
TEXT = RGBColor(245, 247, 252)
MUTED = RGBColor(182, 194, 214)
ACCENT = RGBColor(240, 191, 76)
ACCENT_2 = RGBColor(126, 214, 171)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def panel(slide, left, top, width, height, color=PANEL, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def write_title(slide, title, kicker=None):
    bar = panel(slide, 0, 0, 13.33, 0.82, color=PANEL, radius=False)
    bar.fill.transparency = 0.0
    box = textbox(slide, 0.55, 0.18, 8.2, 0.35)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(25)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if kicker:
        box = textbox(slide, 8.9, 0.19, 3.8, 0.28)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = kicker
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = ACCENT


def footer(slide, text="versebattles.com"):
    box = textbox(slide, 0.55, 7.0, 12.0, 0.2)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_picture(slide, path, left, top, width, height=None):
    if not path.exists():
        return None
    if height is None:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def bullets(slide, items, left, top, width, height, font_size=19):
    box = textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(8)
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT


def caption(slide, title, body, left, top, width, height=0.5):
    panel(slide, left, top, width, height, color=PANEL_2)
    box = textbox(slide, left + 0.14, top + 0.08, width - 0.22, height - 0.1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_2
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = body
    run.font.size = Pt(10.5)
    run.font.color.rgb = TEXT


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_picture(slide, SPLASH_HOPEFUL, 0, 0, 13.33, 7.5)
    overlay = panel(slide, 0, 0, 13.33, 7.5, color=BG, radius=False)
    overlay.fill.transparency = 0.38
    add_picture(slide, LOGO, 0.75, 0.5, 1.25)

    box = textbox(slide, 0.8, 1.45, 7.0, 2.4)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VerseBattles for Christian School Leaders in Uganda"
    run.font.size = Pt(29)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p = tf.add_paragraph()
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = "A browser-based Scripture game that can help schools make Bible memory more active, more repeatable, and easier for students to return to."
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE

    note = panel(slide, 0.82, 4.85, 5.7, 1.2, color=PANEL_2)
    note.fill.transparency = 0.1
    box = textbox(slide, 1.03, 5.05, 5.25, 0.8)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Designed for quick testing on school devices, shared screens, or student phones without requiring an app-store rollout."
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = TEXT

    footer(slide, "Presentation for Christian school leaders")


def side_image_slide(prs, title, items, image, kicker=None, image_caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, title, kicker)
    bullets(slide, items, 0.72, 1.22, 5.0, 5.5)
    panel(slide, 5.92, 1.16, 6.45, 5.55)
    add_picture(slide, image, 6.08, 1.34, 6.1, 4.78)
    if image_caption:
        caption(slide, image_caption[0], image_caption[1], 6.15, 6.02, 5.95)
    footer(slide)


def menu_graphics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, "Easy for Students to Understand", "School-friendly entry points")

    items = [
        "Students do not need a long software lesson before they begin.",
        "The menu gives clear paths such as solo play, missions, learning, and multiplayer.",
        "That makes it easier for a teacher, chaplain, or club leader to explain the activity quickly.",
        "It also gives the school more than one way to use the same tool during the week.",
    ]
    bullets(slide, items, 0.72, 1.18, 4.6, 5.1, font_size=18)

    add_picture(slide, MENU_SCREEN, 5.05, 1.16, 7.55, 3.75)
    caption(slide, "Main menu", "The game already exposes the main activity paths clearly, which reduces first-session confusion.", 5.28, 4.62, 6.9)

    x = 5.22
    for icon in MENU_ICONS:
        panel(slide, x, 5.38, 1.28, 1.1, color=PANEL_2)
        add_picture(slide, icon, x + 0.08, 5.46, 1.12, 0.94)
        x += 1.42

    footer(slide)


def school_use_cases_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, "How a Christian School Could Use It", "Practical use, not theory")

    panel(slide, 0.72, 1.22, 3.85, 4.95)
    panel(slide, 4.74, 1.22, 3.85, 4.95)
    panel(slide, 8.76, 1.22, 3.85, 4.95)

    blocks = [
        ("Bible lesson support", [
            "Use it as a short warm-up before Bible class.",
            "Tie play sessions to the week's memory verse or theme.",
            "Let students revisit Scripture through action, not only recitation."
        ]),
        ("Clubs and houses", [
            "Run short competitions between classes, houses, or Scripture clubs.",
            "Use multiplayer or turn-based shared-device play when equipment is limited.",
            "Give students a reason to come back the next week."
        ]),
        ("Chaplaincy and discipleship", [
            "Use missions, review, and Verse of the Day as follow-up after chapel.",
            "Support students who learn better through repetition and interaction.",
            "Keep Scripture central while still feeling modern to students."
        ]),
    ]

    xs = [0.95, 4.97, 8.99]
    for (head, lines), x in zip(blocks, xs):
        box = textbox(slide, x, 1.55, 3.35, 4.2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = head
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.bullet = True
            p.space_before = Pt(8)
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT

    note = textbox(slide, 0.92, 6.35, 11.5, 0.35)
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The strongest role is usually reinforcement: helping students return to Scripture more often between teaching moments."
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT_2
    footer(slide)


def mosaic_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, "More Than One Kind of Scripture Activity", "Different school rhythms")

    bullets(slide, [
        "Solo play works for individual practice.",
        "Missions add structure for guided progression.",
        "Review mode gives a quieter follow-up after action play.",
        "Multiplayer can support group energy in clubs or supervised sessions.",
        "This gives school leaders flexibility instead of a one-format tool."
    ], 0.72, 1.2, 4.5, 5.3, font_size=18)

    add_picture(slide, MISSIONS, 5.25, 1.2, 3.0, 2.15)
    add_picture(slide, REVIEW, 8.45, 1.2, 3.0, 2.15)
    add_picture(slide, MULTI, 5.25, 3.7, 6.2, 2.25)
    caption(slide, "Missions", "Guided progression is easier to explain than an endless arcade mode.", 5.35, 3.08, 2.75, 0.42)
    caption(slide, "Review", "Students can slow down and focus on the verse after the game pressure.", 8.55, 3.08, 2.75, 0.42)
    caption(slide, "Shared play", "Group play can help clubs or houses turn Bible memory into a repeatable event.", 5.48, 6.08, 5.95, 0.42)
    footer(slide)


def pilot_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, "A Sensible First Pilot", "Start small and judge it honestly")
    add_picture(slide, SPLASH_STRIKE, 7.55, 0.82, 5.78, 6.68)
    overlay = panel(slide, 7.55, 0.82, 5.78, 6.68, color=BG, radius=False)
    overlay.fill.transparency = 0.5

    items = [
        "Test it with one class, one club, or one Scripture-memory group first.",
        "Use the real devices students or staff would actually use.",
        "Watch how long it takes before students understand the activity.",
        "Check whether they come back willingly for a second session.",
        "Decide whether it works best for Bible class, chapel follow-up, clubs, or house competitions.",
    ]
    bullets(slide, items, 0.78, 1.38, 6.1, 4.8, font_size=18)

    quote = panel(slide, 0.82, 5.85, 5.9, 0.88, color=PANEL_2)
    quote.fill.transparency = 0.06
    box = textbox(slide, 1.05, 6.05, 5.45, 0.44)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The goal is not to replace teaching. The goal is to make Scripture engagement easier to repeat."
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT
    footer(slide)


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_picture(slide, SPLASH_HOPEFUL, 0, 0, 13.33, 7.5)
    overlay = panel(slide, 0, 0, 13.33, 7.5, color=BG, radius=False)
    overlay.fill.transparency = 0.3

    add_picture(slide, LOGO, 0.82, 0.58, 1.1)
    card = panel(slide, 0.9, 1.45, 6.2, 4.9, color=PANEL)
    card.fill.transparency = 0.08
    box = textbox(slide, 1.15, 1.75, 5.65, 4.2)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Suggested next step"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    for line in [
        "Open the game and try one short session with a teacher, chaplain, or student leader.",
        "Choose one practical use case: Bible class warm-up, Scripture club, or chapel follow-up.",
        "If it strengthens Scripture memory and student engagement, expand from there.",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    box = textbox(slide, 7.55, 1.95, 4.7, 3.6)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "VerseBattles can be a useful school tool when:"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    for line in [
        "students already have at least some device or shared-screen access",
        "leaders want a more active Scripture-memory rhythm",
        "the school values reinforcement between teaching moments",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
    footer(slide, "VerseBattles for Christian School Leaders in Uganda")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    side_image_slide(
        prs,
        "Why School Leaders May Find It Useful",
        [
            "Many students respond more strongly to active repetition than to one more worksheet.",
            "VerseBattles ties Scripture recall to gameplay, which can make memory work feel more engaging.",
            "Because it is browser-based, a school can test it quickly without planning a full app deployment.",
            "It can support teacher-led Bible learning rather than trying to replace it.",
            "That makes it suitable for a cautious pilot before any broader rollout.",
        ],
        GAMEPLAY,
        kicker="What problem it addresses",
        image_caption=("Scripture through gameplay", "Students answer Bible prompts in order to act, survive, and progress."),
    )
    menu_graphics_slide(prs)
    school_use_cases_slide(prs)
    mosaic_slide(prs)
    side_image_slide(
        prs,
        "Why a Browser-Based Approach Matters",
        [
            "A school can open a URL and test the activity without waiting for app-store installs.",
            "That is useful when devices are shared or when technical setup time is limited.",
            "The same product can be tried in class, after chapel, or in a Scripture club.",
            "School leaders can judge the real student response quickly instead of making a large upfront commitment.",
            "If the first test is weak, the school can stop without being locked into a heavy rollout.",
        ],
        MENU_SCREEN,
        kicker="Low-friction testing",
        image_caption=("Fast to trial", "The product can be assessed as a practical school tool, not only admired in theory."),
    )
    pilot_slide(prs)
    closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
