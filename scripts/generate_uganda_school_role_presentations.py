from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/michael/proj/dcgame")
OUT_DIR = ROOT / "docs/marketing/presentations"

LOGO = ROOT / "images/VerseBattles-logo.png"
SPLASH_HOPEFUL = ROOT / "images/splash_hopeful_clean.png"
SPLASH_STRIKE = ROOT / "images/splash_strike.png"
MENU_SCREEN = ROOT / "public/landing/menu-screen.png"
GAMEPLAY = ROOT / "public/landing/gameplay-screen.png"
MISSIONS = ROOT / "output/web-game/onboarding-missions/shot-0.png"
REVIEW = ROOT / "output/web-game/review-toast-leak/review-no-toast.png"
MULTI = ROOT / "output/web-game/multiplayer-regression/pageA-final.png"

MENU_ICONS = [
    ROOT / "images/menu/solo_game_labeled_icon.png",
    ROOT / "images/menu/missions_labeled_icon.png",
    ROOT / "images/menu/learn_verses_labeled_icon.png",
    ROOT / "images/menu/multiplayer_labeled_icon.png",
    ROOT / "images/menu/groups_labeled_icon.png",
]

BG = RGBColor(14, 20, 31)
PANEL = RGBColor(24, 34, 52)
PANEL_2 = RGBColor(33, 49, 77)
TEXT = RGBColor(245, 247, 252)
MUTED = RGBColor(182, 194, 214)
ACCENT = RGBColor(240, 191, 76)
ACCENT_2 = RGBColor(126, 214, 171)
WHITE = RGBColor(255, 255, 255)


PROFILES = {
    "headteachers": {
        "filename": "VerseBattles_Uganda_Head_Teachers.pptx",
        "audience_footer": "Presentation for head teachers and school administrators",
        "title": "VerseBattles for Head Teachers and School Leaders in Uganda",
        "subtitle": "A browser-based Scripture game that can help a Christian school strengthen Bible-memory routines without a heavy technology rollout.",
        "note": "Designed for small pilots on existing school devices, shared screens, or student phones before any wider decision.",
        "slide2_title": "Why School Leadership May Want to Test It",
        "slide2_kicker": "Operational and educational value",
        "slide2_bullets": [
            "It can give Bible memory a more active format without asking the school to build a large digital programme first.",
            "Because it is browser-based, the school can evaluate it quickly on real devices already in use.",
            "It offers a modest pilot path instead of a costly, school-wide rollout from day one.",
            "It may support student engagement in Scripture while still leaving teaching and discipline in the hands of staff.",
            "That makes it easier to assess honestly before committing time or attention more broadly.",
        ],
        "slide2_caption_title": "A practical pilot candidate",
        "slide2_caption_body": "The key question for leadership is not whether it looks impressive, but whether it works simply and repeatably in the school context.",
        "use_case_title": "Where It Could Fit in School Life",
        "use_case_kicker": "Administrative use cases",
        "use_blocks": [
            ("Bible and RE periods", [
                "Use it as a short warm-up or reinforcement activity.",
                "Tie sessions to the week's Scripture theme or memory verse.",
                "Support retention without replacing teacher instruction.",
            ]),
            ("Clubs and houses", [
                "Use it for Scripture clubs, inter-house activities, or supervised competitions.",
                "Create repeatable student interest around Bible knowledge.",
                "Let the same tool work in more than one school setting.",
            ]),
            ("Chapel follow-up", [
                "Use it after chapel to help students revisit the theme during the week.",
                "Give chaplains and teachers a lighter-prep follow-up option.",
                "Turn one sermon moment into repeated Scripture contact.",
            ]),
        ],
        "browser_title": "Why the Browser-Based Model Matters",
        "browser_kicker": "Lower rollout friction",
        "browser_bullets": [
            "A school can test it without planning an app installation process for every device.",
            "That matters when devices are shared, student access is uneven, or staff time is limited.",
            "Leadership can assess the real value first, then decide whether broader use is justified.",
            "If the pilot is weak, the school can stop early without sunk rollout complexity.",
            "If the pilot is strong, expansion can happen in stages rather than all at once.",
        ],
        "browser_caption_title": "Fast to evaluate",
        "browser_caption_body": "This is easier to test as a school tool than a platform that requires a long technical setup first.",
        "pilot_title": "A Responsible First Pilot",
        "pilot_kicker": "Start with evidence, not assumptions",
        "pilot_bullets": [
            "Choose one class, club, or chapel follow-up group first.",
            "Use the real devices students and staff would actually rely on.",
            "Watch how much explanation is required before meaningful use begins.",
            "Check whether students return for a second session willingly.",
            "Decide whether the best fit is class support, clubs, house competitions, or chaplaincy follow-up.",
        ],
        "pilot_quote": "The most important leadership question is simple: does this strengthen Scripture engagement enough to justify school attention?",
        "closing_title": "Suggested next step",
        "closing_bullets": [
            "Ask one trusted teacher or chaplain to run a small trial.",
            "Choose one defined use case, not many at once.",
            "Review whether it improves Scripture engagement, repeatability, and ease of use.",
        ],
        "closing_side_title": "A good leadership fit usually means:",
        "closing_side_bullets": [
            "the school wants stronger Scripture reinforcement during the week",
            "shared-device or browser access is at least somewhat realistic",
            "leaders prefer a cautious pilot before wider adoption",
        ],
    },
    "chaplains": {
        "filename": "VerseBattles_Uganda_Chaplains_Bible_Teachers.pptx",
        "audience_footer": "Presentation for chaplains and Bible teachers",
        "title": "VerseBattles for Chaplains and Bible Teachers in Uganda",
        "subtitle": "A browser-based Scripture game that can help students return to Bible memory with more energy, repetition, and attention.",
        "note": "Designed to support teacher-led discipleship, Bible teaching, and chapel follow-up rather than replace them.",
        "slide2_title": "Why Chaplains and Bible Teachers May Find It Valuable",
        "slide2_kicker": "Teaching and discipleship value",
        "slide2_bullets": [
            "Many students need more than recitation alone if Scripture memory is going to stay alive during the week.",
            "VerseBattles ties Bible recall to action, which can make repetition feel less passive.",
            "It gives students another reason to come back to the verse after the first lesson or chapel message.",
            "The game still keeps Scripture central rather than using Bible content as a thin decoration.",
            "That makes it useful as reinforcement around teaching, not only as entertainment.",
        ],
        "slide2_caption_title": "Scripture through action",
        "slide2_caption_body": "Students answer Bible prompts in order to move, attack, and progress, which turns recall into a more active exercise.",
        "use_case_title": "How It Could Support Ministry in School",
        "use_case_kicker": "Chaplains and Bible teachers",
        "use_blocks": [
            ("Bible class reinforcement", [
                "Use it before or after class to revisit the week's passage.",
                "Help students keep returning to Scripture between lessons.",
                "Support different learning styles through repeated interaction.",
            ]),
            ("Chapel follow-up", [
                "Connect the week's chapel message to a Scripture-memory rhythm.",
                "Let students revisit the passage through missions, review, or learning modes.",
                "Extend the ministry value of chapel beyond one gathering.",
            ]),
            ("Student discipleship groups", [
                "Use it in Scripture clubs, discipleship groups, or house fellowships.",
                "Create a shared activity that still keeps the Bible content central.",
                "Give students a reason to re-engage with verses together.",
            ]),
        ],
        "browser_title": "Why It Can Be Easier to Introduce",
        "browser_kicker": "Less friction for ministry use",
        "browser_bullets": [
            "A chaplain or teacher can test it quickly without waiting for a full app rollout.",
            "That makes it easier to try during a real week of school life rather than in a special technical setup.",
            "The same tool can be used in Bible class, after chapel, or in a student discipleship group.",
            "If students respond well, it can become part of a wider Scripture-memory rhythm.",
            "If they do not, the school has not overcommitted resources to find that out.",
        ],
        "browser_caption_title": "Simple to trial",
        "browser_caption_body": "For ministry staff, low-friction testing matters because the best tool is the one that can actually be used in a normal school week.",
        "pilot_title": "A Strong First Ministry Pilot",
        "pilot_kicker": "Test it in a real discipleship rhythm",
        "pilot_bullets": [
            "Choose one verse theme, one class, or one discipleship group first.",
            "Use it with the same devices students would really have access to.",
            "Check whether students understand the activity quickly enough to use it without heavy explanation every time.",
            "Watch whether they return for a second or third session with interest.",
            "Decide whether the strongest fit is Bible class, chapel follow-up, or Scripture club use.",
        ],
        "pilot_quote": "The right question is not only, 'Did students enjoy it?' but also, 'Did it help them return to Scripture more willingly and more often?'",
        "closing_title": "Suggested next step",
        "closing_bullets": [
            "Run one short test with a class, Bible teacher, or chaplaincy group.",
            "Tie the pilot to a real passage already being taught.",
            "Judge whether it helps students return to Scripture with more attention and consistency.",
        ],
        "closing_side_title": "A strong ministry fit usually means:",
        "closing_side_bullets": [
            "leaders want a more active Scripture-memory rhythm",
            "students need reinforcement between teaching moments",
            "the school values tools that support, rather than replace, teacher-led discipleship",
        ],
    },
    "owners": {
        "filename": "VerseBattles_Uganda_School_Owners_Board_Sponsors.pptx",
        "audience_footer": "Presentation for school owners, board members, and sponsors",
        "title": "VerseBattles for School Owners, Board Members, and Sponsors in Uganda",
        "subtitle": "A browser-based Scripture game that can help a Christian school test a stronger Bible-memory rhythm without taking on the risk of a heavy technology project.",
        "note": "Best evaluated as a small, practical pilot first, using the devices and school conditions that already exist on the ground.",
        "slide2_title": "Why Decision-Makers May Want to Look at It",
        "slide2_kicker": "Stewardship and school value",
        "slide2_bullets": [
            "It offers a modest entry point into digital Scripture engagement without forcing a large procurement decision first.",
            "Because it is browser-based, the school can test the real educational and ministry value before making wider commitments.",
            "That lowers the risk of spending time or money on a tool that looks promising in theory but proves awkward in practice.",
            "If it works well, it can strengthen the school's Christian identity in a visible, student-facing way.",
            "If it does not work well, the school can stop early with limited operational cost.",
        ],
        "slide2_caption_title": "A cautious investment posture",
        "slide2_caption_body": "The right decision is not based on novelty. It is based on whether the tool creates real Scripture engagement under normal school conditions.",
        "use_case_title": "Where It Could Add Value to the School",
        "use_case_kicker": "Institution-level use cases",
        "use_blocks": [
            ("Christian distinctiveness", [
                "Support the school's claim to provide more than academic instruction alone.",
                "Give students a modern Scripture activity that still keeps the Bible central.",
                "Strengthen the visible Christian culture of the school.",
            ]),
            ("Student engagement", [
                "Create a repeatable Bible-memory rhythm that students may return to more willingly.",
                "Give Bible learning a format that feels more active than routine recitation alone.",
                "Support clubs, classes, and chaplaincy programmes with one shared tool.",
            ]),
            ("Pilot-friendly innovation", [
                "Try one contained pilot instead of launching a major digital initiative.",
                "Judge results before wider rollout or sponsorship.",
                "Expand only if the value is clear enough to justify further support.",
            ]),
        ],
        "browser_title": "Why the Delivery Model Matters to Sponsors and Boards",
        "browser_kicker": "Lower commitment before proof",
        "browser_bullets": [
            "A board or sponsor does not have to underwrite a large technical rollout just to see whether the concept works.",
            "The school can assess the tool under ordinary conditions with ordinary devices first.",
            "That allows more disciplined stewardship than approving a larger programme before the evidence is visible.",
            "If the pilot is effective, support can then be directed toward the next practical step rather than speculation.",
            "This makes it easier to fund measured progress instead of vague digital ambition.",
        ],
        "browser_caption_title": "Proof before scale",
        "browser_caption_body": "A sponsor or board can ask for visible pilot evidence before considering any broader backing.",
        "pilot_title": "What a Wise First Pilot Looks Like",
        "pilot_kicker": "Small, observable, accountable",
        "pilot_bullets": [
            "Run the pilot with one class, one club, or one chapel follow-up group first.",
            "Use the devices students and staff would actually rely on if the tool continued.",
            "Watch startup friction, student understanding, and whether the second session still has value.",
            "Ask teachers and chaplains whether it genuinely strengthens Scripture engagement or only creates short-term novelty.",
            "Make any next-stage decision only after that evidence is visible.",
        ],
        "pilot_quote": "For boards and sponsors, the key issue is stewardship: does this create enough real ministry and educational value to deserve further support?",
        "closing_title": "Suggested next step",
        "closing_bullets": [
            "Approve or support one limited pilot rather than a broad rollout.",
            "Ask for clear feedback from school leadership, chaplains, and teachers.",
            "If the school sees genuine Scripture-memory value, fund the next practical stage from evidence, not assumption.",
        ],
        "closing_side_title": "A strong ownership or sponsor fit usually means:",
        "closing_side_bullets": [
            "the school wants visible Christian distinctiveness in student life",
            "leaders prefer measured pilots before larger commitments",
            "supporters want practical ministry value rather than technology theatre",
        ],
    },
}


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def panel(slide, left, top, width, height, color=PANEL, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def write_title(slide, title, kicker=None):
    panel(slide, 0, 0, 13.33, 0.82, color=PANEL, radius=False)
    box = textbox(slide, 0.55, 0.18, 8.2, 0.35)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(25)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if kicker:
        box = textbox(slide, 8.9, 0.19, 3.8, 0.28)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = kicker
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = ACCENT


def footer(slide, text="versebattles.com"):
    box = textbox(slide, 0.55, 7.0, 12.0, 0.2)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_picture(slide, path, left, top, width, height=None):
    if not path.exists():
        return None
    if height is None:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def bullets(slide, items, left, top, width, height, font_size=19):
    box = textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(8)
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT


def caption(slide, title, body, left, top, width, height=0.5):
    panel(slide, left, top, width, height, color=PANEL_2)
    box = textbox(slide, left + 0.14, top + 0.08, width - 0.22, height - 0.1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_2
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = body
    run.font.size = Pt(10.5)
    run.font.color.rgb = TEXT


def title_slide(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_picture(slide, SPLASH_HOPEFUL, 0, 0, 13.33, 7.5)
    overlay = panel(slide, 0, 0, 13.33, 7.5, color=BG, radius=False)
    overlay.fill.transparency = 0.38
    add_picture(slide, LOGO, 0.75, 0.5, 1.25)

    box = textbox(slide, 0.8, 1.45, 7.0, 2.4)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg["title"]
    run.font.size = Pt(29)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p = tf.add_paragraph()
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = cfg["subtitle"]
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE

    note = panel(slide, 0.82, 4.85, 5.8, 1.2, color=PANEL_2)
    note.fill.transparency = 0.1
    box = textbox(slide, 1.03, 5.05, 5.35, 0.8)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg["note"]
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = TEXT
    footer(slide, cfg["audience_footer"])


def side_image_slide(prs, title, items, image, kicker=None, image_caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, title, kicker)
    bullets(slide, items, 0.72, 1.22, 5.0, 5.5)
    panel(slide, 5.92, 1.16, 6.45, 5.55)
    add_picture(slide, image, 6.08, 1.34, 6.1, 4.78)
    if image_caption:
        caption(slide, image_caption[0], image_caption[1], 6.15, 6.02, 5.95)
    footer(slide)


def menu_graphics_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, title, subtitle)

    items = [
        "Students do not need a long software lesson before they begin.",
        "The menu gives clear paths such as solo play, missions, learning, and multiplayer.",
        "That makes it easier for staff to explain the activity quickly and supervise it realistically.",
        "It also gives the school more than one way to use the same tool during the week.",
    ]
    bullets(slide, items, 0.72, 1.18, 4.6, 5.1, font_size=18)

    add_picture(slide, MENU_SCREEN, 5.05, 1.16, 7.55, 3.75)
    caption(slide, "Main menu", "The product surface is clear enough to support a real pilot rather than only a polished demo.", 5.28, 4.62, 6.9)

    x = 5.22
    for icon in MENU_ICONS:
        panel(slide, x, 5.38, 1.28, 1.1, color=PANEL_2)
        add_picture(slide, icon, x + 0.08, 5.46, 1.12, 0.94)
        x += 1.42

    footer(slide)


def use_cases_slide(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, cfg["use_case_title"], cfg["use_case_kicker"])

    panel(slide, 0.72, 1.22, 3.85, 4.95)
    panel(slide, 4.74, 1.22, 3.85, 4.95)
    panel(slide, 8.76, 1.22, 3.85, 4.95)

    xs = [0.95, 4.97, 8.99]
    for (head, lines), x in zip(cfg["use_blocks"], xs):
        box = textbox(slide, x, 1.55, 3.35, 4.2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = head
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.bullet = True
            p.space_before = Pt(8)
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT

    note = textbox(slide, 0.92, 6.35, 11.5, 0.35)
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The strongest role is usually reinforcement: helping students return to Scripture more often between teaching moments."
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT_2
    footer(slide)


def mosaic_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, "More Than One Kind of Scripture Activity", "Different school rhythms")

    bullets(slide, [
        "Solo play works for individual practice.",
        "Missions add structure for guided progression.",
        "Review mode gives a quieter follow-up after action play.",
        "Multiplayer can support group energy in clubs or supervised sessions.",
        "This gives school leaders flexibility instead of a one-format tool."
    ], 0.72, 1.2, 4.5, 5.3, font_size=18)

    add_picture(slide, MISSIONS, 5.25, 1.2, 3.0, 2.15)
    add_picture(slide, REVIEW, 8.45, 1.2, 3.0, 2.15)
    add_picture(slide, MULTI, 5.25, 3.7, 6.2, 2.25)
    caption(slide, "Missions", "Guided progression is easier to explain than an endless arcade mode.", 5.35, 3.08, 2.75, 0.42)
    caption(slide, "Review", "Students can slow down and focus on the verse after the game pressure.", 8.55, 3.08, 2.75, 0.42)
    caption(slide, "Shared play", "Group play can help clubs or houses turn Bible memory into a repeatable event.", 5.48, 6.08, 5.95, 0.42)
    footer(slide)


def browser_slide(prs, cfg):
    side_image_slide(
        prs,
        cfg["browser_title"],
        cfg["browser_bullets"],
        MENU_SCREEN,
        kicker=cfg["browser_kicker"],
        image_caption=(cfg["browser_caption_title"], cfg["browser_caption_body"]),
    )


def pilot_slide(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    write_title(slide, cfg["pilot_title"], cfg["pilot_kicker"])
    add_picture(slide, SPLASH_STRIKE, 7.55, 0.82, 5.78, 6.68)
    overlay = panel(slide, 7.55, 0.82, 5.78, 6.68, color=BG, radius=False)
    overlay.fill.transparency = 0.5

    bullets(slide, cfg["pilot_bullets"], 0.78, 1.38, 6.1, 4.8, font_size=18)

    quote = panel(slide, 0.82, 5.85, 5.95, 0.88, color=PANEL_2)
    quote.fill.transparency = 0.06
    box = textbox(slide, 1.05, 6.05, 5.45, 0.44)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg["pilot_quote"]
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = TEXT
    footer(slide)


def closing_slide(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_picture(slide, SPLASH_HOPEFUL, 0, 0, 13.33, 7.5)
    overlay = panel(slide, 0, 0, 13.33, 7.5, color=BG, radius=False)
    overlay.fill.transparency = 0.3

    add_picture(slide, LOGO, 0.82, 0.58, 1.1)
    card = panel(slide, 0.9, 1.45, 6.2, 4.9, color=PANEL)
    card.fill.transparency = 0.08
    box = textbox(slide, 1.15, 1.75, 5.65, 4.2)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg["closing_title"]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    for line in cfg["closing_bullets"]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    box = textbox(slide, 7.55, 1.95, 4.7, 3.6)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg["closing_side_title"]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    for line in cfg["closing_side_bullets"]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_before = Pt(10)
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
    footer(slide, cfg["title"])


def build_profile(name, cfg):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs, cfg)
    side_image_slide(
        prs,
        cfg["slide2_title"],
        cfg["slide2_bullets"],
        GAMEPLAY,
        kicker=cfg["slide2_kicker"],
        image_caption=(cfg["slide2_caption_title"], cfg["slide2_caption_body"]),
    )
    menu_graphics_slide(prs, "Easy for Students to Understand", "School-friendly entry points")
    use_cases_slide(prs, cfg)
    mosaic_slide(prs)
    browser_slide(prs, cfg)
    pilot_slide(prs, cfg)
    closing_slide(prs, cfg)

    output = OUT_DIR / cfg["filename"]
    prs.save(output)
    print(f"Wrote {output}")


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for name, cfg in PROFILES.items():
        build_profile(name, cfg)


if __name__ == "__main__":
    main()
